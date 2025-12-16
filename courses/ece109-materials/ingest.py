from pathlib import Path
from pptx import Presentation
import subprocess

RAW = Path("raw")
TEXT = Path("text")
TEXT.mkdir(exist_ok=True)

def pdf_to_text(pdf):
    out = TEXT / f"{pdf.stem}.txt"
    subprocess.run(["pdftotext", pdf, out])

def pptx_to_text(pptx):
    out = TEXT / f"{pptx.stem}.txt"
    prs = Presentation(pptx)
    with open(out, "w") as f:
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    f.write(shape.text + "\n")

for file in RAW.iterdir():
    if file.suffix == ".pdf":
        pdf_to_text(file)
    elif file.suffix == ".pptx":
        pptx_to_text(file)

print("Documents converted to text.")
